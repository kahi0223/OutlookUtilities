from pptx import Presentation
from pptx.util import Pt

def make_slide(name_layout, name_title):
    """ create new slide
    """
    slide = Presentation.slides.add_slide(name_layout)
    if name_title:
        title = slide.shapes.title
        title.text = name_title
    return slide

def input_text(slide, shape_number, input_dic):
    """input text in placeholder
    - decide which placeholoder by shape number
    - input text and new line by dict 
    """
    for n, t in input_dic.items():
        slide.shapes[shape_number].text_frame.add_paragraph()
        text_box = slide.shapes[shape_number].text_frame.paragraphs[n]
        text_box.text = t

def change_font(slide, shape_number):
    """can change font by lines
    """
    for i in [0,9]:
        slide.shapes[shape_number].text_frame.paragraphs[i].font.underline = True #under line
    for i in [1,3,5,7]:
        slide.shapes[shape_number].text_frame.paragraphs[i].level = 1  # one indent
    for i in [2,4,6]:
        slide.shapes[shape_number].text_frame.paragraphs[i].level = 2  # two indents

if __name__ == "__main__":

    PRESENTATION_NAME = "template.pptx"

    Presentation = Presentation(PRESENTATION_NAME)

    # setting layout (need to make master layout before run)
    TITLE_LAYOUT = Presentation.slide_layouts[0]
    BASIC_LAYOUT = Presentation.slide_layouts[1]
    DOUBLE_SHAPES_LAYOUT = Presentation.slide_layouts[2]
    END_LAYOUT = Presentation.slide_layouts[3]

    # create title slide
    presentation_title = "Presentation Title"
    make_slide(TITLE_LAYOUT, presentation_title)

    # create normal slide
    slide1_title = "slide1 title"
    slide1 = make_slide(BASIC_LAYOUT, slide1_title)

    input_text_dict1 ={
        0: "first line",
        1: "second line",
        2: "third line",
        3: "forth line",
        4: "fifth line"
    }
    input_text(slide1, 1, input_text_dict1)

    # create slide and change font
    slide2_title = "slide2 title"
    input_text_dict2 ={
        0: "first line",
        1: "second line",
        2: "third line",
        3: "forth line",
        4: "fifth line",
        5: "sixth line",
        6: "seventh line",
        7: "eighth line",
        8: "ninth line",
        9: "tenth line"
    }
    slide2 = make_slide(BASIC_LAYOUT, slide2_title)
    input_text(slide2, 1, input_text_dict2)
    change_font(slide2, 1)

    # create double slide
    slide3_title = "slide3 title"
    input_text_dict3 ={
        0: "first line",
        1: "second line",
        2: "third line",
        3: "forth line",
        4: "fifth line"
    }

    input_text_dict4 ={
        0: "first line",
        1: "second line",
        2: "third line",
        3: "forth line",
        4: "fifth line"
    }
    slide3 = make_slide(DOUBLE_SHAPES_LAYOUT, slide3_title)
    input_text(slide3, 1, input_text_dict3)
    input_text(slide3, 2, input_text_dict4)

    # create end slide
    make_slide(END_LAYOUT, None)

    Presentation.save(PRESENTATION_NAME)
    print("Saved")