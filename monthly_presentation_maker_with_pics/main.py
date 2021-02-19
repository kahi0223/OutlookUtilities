from pptx import Presentation
from pptx.util import Pt
from PIL import Image
import glob

#change here
this_month = "MM"
pre_month = "pre MM"
year = "YYYY"

presentation_name = "report_ppt.pptx"
my_name = "NAME"

source_pic_name = "raw_pic"
source_pic_extention = ".jpg"

picts = {
    "shaped_pic1":source_pic_name + "1" + source_pic_extention,
    "shaped_pic2":source_pic_name + "2" + source_pic_extention}


# スライドの生成関数
def make_slide(name_layout, name_title):
    slide = presentation.slides.add_slide(name_layout)
    title = slide.shapes.title
    title.text = name_title
    return slide

# テキストの入力関数
def input_text(slide, shape_number, input_dic):
    slide.shapes[shape_number].text_frame
    slide.shapes[shape_number].text_frame.clear()    

    for n, t in input_dic.items():
        slide.shapes[shape_number].text_frame.add_paragraph()     
        text_box = slide.shapes[shape_number].text_frame.paragraphs[n]
        text_box.text = t
  

# テキストボックス数を数える補助関数(未使用)
def count_text_box(slide): 
    for shape in slide.placeholders:
        print(shape.placeholder_format.idx, shape.name)

# インデントを下げる関数
def add_indent(slide, shape_number, indent_number, start, end):
    for i in list(range(start,end+1)):
        slide.shapes[shape_number].text_frame.paragraphs[i].level = indent_number


# アンダーライン関数
def under_line(slide, shape_number, paragraph_number):
    slide.shapes[shape_number].text_frame.paragraphs[paragraph_number].font.underline = True

# 画像挿入関数(左右中心・下揃え)
def _add_pict(slide, pic_pass, times):
    picture = slide.shapes.add_picture(pic_pass,0,0)
    base = picture.height 
    picture.height = int(presentation.slide_height * times)
    picture.width = int(picture.width * (picture.height / base))
    picture.left = int((presentation.slide_width - picture.width) / 2)
    picture.top = int((presentation.slide_height - picture.height) - base/20)
    return picture

# 画像挿入関数(左右中心・高さ調整)
def _add_pict2(slide, pic_pass, times):
    picture = slide.shapes.add_picture(pic_pass,0,0)
    base = picture.height 
    picture.height = int(presentation.slide_height * times)
    picture.width = int(picture.width * (picture.height / base))
    picture.left = int((presentation.slide_width - picture.width) / 2)
    picture.top = int((presentation.slide_height - picture.height) - base/5) # change here
    return picture


#  画像をつけたスライド作成関数・下揃え
def add_pict_slide(title, text_dic, pict):
    slide = make_slide(PICT_LAYOUT, title)
    shape = 1
    input_text(slide, shape, text_dic)
    under_line(slide, shape, 0)
    add_indent(slide, shape, 1, 1, 3)
    _add_pict(slide, pict, 0.5)  # change here
    return slide

#  画像をつけたスライド作成関数・好きな場所
def add_pict2_slide(title, text_dic, pict):
    slide = make_slide(PICT_LAYOUT, title)
    shape = 1
    input_text(slide, shape, text_dic)
    under_line(slide, shape, 0)
    add_indent(slide, shape, 1, 1, 3)
    _add_pict2(slide, pict, 0.3)  # change here
    return slide

if __name__== "__main__":

    # 画像切り取り

    for name, image_path in picts.items():
        image = Image.open(image_path)
        new_name = name + ".jpg"
        crapped_image = image.crop(box = (100,200,400,300)) # change here
        crapped_image.save(new_name,format = "jpeg")

    # スライドlayoutの定義
    presentation = Presentation(presentation_name)
    TITLE_LAYOUT = presentation.slide_layouts[0]
    SUMMARY_LAYOUT = presentation.slide_layouts[1]
    PICT_LAYOUT = presentation.slide_layouts[2]


    slide = make_slide(TITLE_LAYOUT, "{} Summary Report".format(this_month))
    input_text(slide, 1, {0: my_name})


    # summary slide
    slide = make_slide(SUMMARY_LAYOUT, this_month +" Summary")

    summary_comment = {
        0: "Summary",
        1: "{}: ##".format(this_month),
        2: "{}: ##".format(pre_month)
        }
    shape = 1
    input_text(slide, shape, summary_comment)

    summary_comment2 = {
        0: "Summary2",
        1: "{}: ##".format(this_month),
        2: "{}: ##".format(pre_month)
        }
    shape = 2
    input_text(slide, shape, summary_comment2)


    ## 画像付きスライドを生成
    text_dict = {
        0: this_month + " (vs " + pre_month + ")",
        1: "memo1",
        2: "memo2",
        3: "memo3"
        }
    add_pict_slide("TITLE1", text_dict, "shaped_pic1.jpg")

    text_dict = {
        0: this_month + " (vs " + pre_month + ")",
        1: "memo4",
        2: "memo5",
        3: "memo6"
        }
    add_pict2_slide("TITLE2", text_dict, "shaped_pic2.jpg")

    presentation.save(presentation_name)